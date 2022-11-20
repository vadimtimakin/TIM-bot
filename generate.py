import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from generate_plot import generate_plot

import json
from docx import Document
import re
from docx.shared import Inches

base_json = json.load(open("data.json", "r"))

def pad(text):
    while len(text) % 8 != 0:
        text += b' '
    return text


def clear_answers(json_answers):
    for key in json_answers.keys():
        if '.jpg' not in json_answers[key] or '.png' not in json_answers[key]:
            json_answers[key] = json_answers[key].strip().strip('.').strip()
    special = ['{{Income}}', '{{NetProfit}}', '{{EBITDA}}']
    for key in special:
        if key in json_answers.keys():
            values = json_answers[key].split('\n')
            for i in range(len(values)):
                v = values[i]
                if not v.isdigit():
                    v = float(re.findall('[\d]+[.,\d]+|[\d]*[.][\d]+|[\d]+', v)[0])
                    if v == int(v):
                        v = int(v)
                    values[i] = v
                v = str(v)
                values[i] = v
                new_key = '{{' + key[2:-2].lower() + str(2020+i) + '}}'
                json_answers[new_key] = values[i]
    if '{{MembersDesc}}' in json_answers.keys():
        text = json_answers['{{MembersDesc}}']
        lines = text.split('\n')
        if len(lines) > 5: lines = lines[:5]
        for line in lines:
            try:
                fio, job, experience = line.split(' ')
            except:
                json_answers['fio'+str(lines.index(line))] = ''
                json_answers['job'+str(lines.index(line))] = ''
                json_answers['experience'+str(lines.index(line))] = ''

            json_answers['fio'+str(lines.index(line))] = fio
            json_answers['job'+str(lines.index(line))] = job
            json_answers['experience'+str(lines.index(line))] = experience
        for i in range(5 - len(lines)):
            json_answers['fio'+str(lines.index(line))] = ''
            json_answers['job'+str(lines.index(line))] = ''
            json_answers['experience'+str(lines.index(line))] = ''

    for x in ['{{memberphoto' + str(i)+'}}' for i in range(1, 6)]:
        json_answers[x] = ''
    if 'path' in json_answers.keys():
        PATH = json_answers['path']
        json_answers['{{CompanyLogo}}'] = PATH + '{{CompanyLogo}}_1.jpg'
        json_answers['{{ProductPhoto}}'] = PATH + '{{ProductPhoto}}_1.jpg'
        for i in range(1, 6):
            cur_key = '{{memberphoto' + str(i) + '}}'
            json_answers[cur_key] = PATH + '{{MembersPhoto}}_' + str(i) + '.jpg'
    news = ['{{new' + str(i) + '}}' for i in range(1, 16)]
    f = 1
    for key in news:
        if key in json_answers.keys():
            f = 0
        else:
            json_answers[key] = ''
    if f:
        json_answers['{{new1}}'] = 'Нет дополнительных вопросов'
    return json_answers   

def generate_file(json_answers):
    doc = Document("MEMO_templateUPD.docx") #6 7
    photo_keys = ['{{CompanyLogo}}', '{{ProductPhoto}}', '{{MembersPhoto}}'] + ['{{memberphoto' + str(i)+'}}' for i in range(1, 6)]
    for i in range(len(doc.paragraphs)):
        p = doc.paragraphs[i]
        for j in range(len(p.runs)):
            cur_phrase = p.runs[j].text
            if cur_phrase in json_answers.keys():
                if cur_phrase not in photo_keys:
                    p.runs[j].text = json_answers[cur_phrase]
                else:
                    run = p.runs[j]
                    run.text = ''
                    try:
                        run.add_picture(json_answers[cur_phrase], height = Inches(2.5))
                    except:
                        pass
    doc.save(f'{json_answers["path"]}MEMOresult.docx')
    
def generate_canvas(json_answers):
    doc = Document("CANVAS_template.docx") #6 7
    photo_keys = ['{{CompanyLogo}}', '{{ProductPhoto}}', '{{MembersPhoto}}'] + ['{{memberphoto' + str(i)+'}}' for i in range(1, 6)]
    for i in range(len(doc.paragraphs)):
        p = doc.paragraphs[i]
        for j in range(len(p.runs)):
            cur_phrase = p.runs[j].text
            if cur_phrase in json_answers.keys():
                if cur_phrase not in photo_keys:
                    p.runs[j].text = json_answers[cur_phrase]
                else:
                    run = p.runs[j]
                    run.text = ''
                    try:
                        run.add_picture(json_answers[cur_phrase], height = Inches(2.5))
                    except:
                        pass
    doc.save(f'{json_answers["path"]}CANVASresult.docx')


def translate(json_file, output_json):
    input_json = {}
    for key in json_file.keys():
        x = key.replace('{', '')
        x = x.replace('}', '')
        input_json[x] = json_file[key]

    output_json['first_slide']['image_path'] = input_json["path"] + "{{ProductPhoto}}_1.jpg"
    output_json['first_slide']['title'] = input_json['ProductName']
    output_json['first_slide']['subtitle'] = input_json['ProductFunc']

    output_json['second_slide']["text_1_2"] = input_json['MarketName']
    output_json['second_slide']["text_2_1"] = input_json['MarketDesc']
    output_json['second_slide']["text_2_2"] = input_json['ProductDesc']

    output_json['third_slide']["text_1"] = input_json['CompetitorsDesc']
    output_json['third_slide']["text_2"] = input_json['ProductUnic']
    output_json['third_slide']["text_3"] = input_json['ProductBest']

    output_json['fifth_slide']["text_1"] = input_json['InvestmentPlans']
    output_json['fifth_slide']["text_2"] = input_json['InvestmentGoal']

    for i in range(1, 6):
        if not f'fio{i}' in input_json:
            input_json[f'fio{i}'] = ""

    for i in range(1, 6):
        if not f'job{i}' in input_json:
            input_json[f'job{i}'] = ""
    
    for i in range(1, 6):
        if not f'experience{i}' in input_json:
            input_json[f'experience{i}'] = ""

    output_json["sixth_slide"]["members"] = {

    input_json['fio1'] : {
            "position": input_json["job1"],
            "experience" : input_json["experience1"]
        },
    input_json['fio2'] : {
            "position": input_json["job2"],
            "experience" : input_json["experience3"]
        },
    input_json['fio3'] : {
            "position": input_json["job3"],
            "experience" : input_json["experience3"]
        },
    input_json['fio4'] : {
            "position": input_json["job4"],
            "experience" : input_json["experience4"]
        },
    input_json['fio5'] : {
            "position": input_json["job5"],
            "experience" : input_json["experience5"]
        },
    }

    output_json["seventh_slide"]["name"] = input_json["AuthorName"]
    output_json["seventh_slide"]["number"] = input_json["AuthorPhone"]
    output_json["seventh_slide"]["email"] = input_json["AuthorEmail"]
    output_json["seventh_slide"]["telegram"] = input_json["AuthorTelegram"]

    return output_json


def generate_pptx(input_json):
    data = translate(input_json, base_json)
    pres = Presentation('slide_plan.pptx')

    def first():
        slide = pres.slides[0]

        # title
        slide.shapes.title.text = data['first_slide']['image_path']
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # subtitle
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left * 2, top * 1.5, width * 10, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = data["first_slide"]["subtitle"]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)

        # image
        top_img_path = data['first_slide']['image_path']

        top_pic = slide.shapes.add_picture(
            top_img_path,
            left=Inches(2),
            top=Inches(3.3),
            width=Inches(6)
        )

        ref_element = slide.shapes[0]._element
        ref_element.addnext(top_pic._element)


    def second():
        slide = pres.slides[1]

        slide.shapes[0].text_frame.text = data["second_slide"]["text_1_1"]
        slide.shapes[1].text_frame.text = data["second_slide"]["text_1_2"]

        slide.shapes[2].text_frame.text = data["second_slide"]["text_2_1"]
        slide.shapes[1].text_frame.paragraphs[0].font.size = Pt(16)

        slide.shapes[3].text_frame.text = data["second_slide"]["text_2_2"]
        slide.shapes[3].text_frame.paragraphs[0].font.size = Pt(16)


    def third():
        slide = pres.slides[2]
        # title
        slide.shapes[1].text = data["third_slide"]["title"]
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        # text
        slide.shapes[4].text = data["third_slide"]["paragraph_1"] + data["third_slide"]["text_1"] +\
                            data["third_slide"]["paragraph_2"] + data["third_slide"]["text_2"] +\
                            data["third_slide"]["paragraph_3"] + data["third_slide"]["text_3"]
        for paragraph in slide.shapes[4].text_frame.paragraphs:
            paragraph.font.size = Pt(16)

    def fourth():

        def add_paragraph(text):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)

        slide = pres.slides[3]
        next_year, current_year = sorted(data["fourth_slide"]["years"], key=int, reverse=True)[:2]

        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left*0.2, top * 3, width*2, height*4)
        tf = txBox.text_frame
        # выручка текущего года
        add_paragraph(f"{data['fourth_slide']['paragraph_1']}({current_year}) - {data['fourth_slide']['years'][str(current_year)][data['fourth_slide']['paragraph_1']]}")
        # прибыль текущего года
        add_paragraph(f"{data['fourth_slide']['paragraph_2']}({current_year}) - {data['fourth_slide']['years'][str(current_year)][data['fourth_slide']['paragraph_2']]}")
        # выручка следующего года
        add_paragraph(f"{data['fourth_slide']['paragraph_3']}({next_year}) - {data['fourth_slide']['years'][str(next_year)][data['fourth_slide']['paragraph_1']]}")
        # прибыль следующего года
        add_paragraph(f"{data['fourth_slide']['paragraph_4']}({next_year}) - {data['fourth_slide']['years'][str(next_year)][data['fourth_slide']['paragraph_2']]}")

        # plot
        generate_plot()
        top_pic = slide.shapes.add_picture(
            data["fourth_slide"]["path2plot"],
            left=Inches(4),
            top=Inches(2),
            width=Inches(6)
        )

        ref_element = slide.shapes[0]._element
        ref_element.addnext(top_pic._element)

    def fifth():
        slide = pres.slides[4]

        # first text block
        slide.shapes[3].text = data["fifth_slide"]["text_1"]
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # second text block
        slide.shapes[4].text = data["fifth_slide"]["text_2"]
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    def sixth():
        slide = pres.slides[5]
        left = top = width = height = Inches(1)
        txBox = slide.shapes.add_textbox(left * 2, top * 1.7, width * 10, height)
        tf = txBox.text_frame

        for name, mem_data in data["sixth_slide"]["members"].items():
            p = tf.add_paragraph()
            p.text = f"• {name}\n\tposition: {mem_data['position']}\n\texperience: {mem_data['experience']}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0,0,0)

    def seventh():
        slide = pres.slides[6]
        slide.shapes[4].text = f"ФИО: {data['seventh_slide']['name']}\n" +\
                                f"Номер телефона: {data['seventh_slide']['number']}\n" + \
                            f"Электронная почта: {data['seventh_slide']['email']}\n" + \
                            f"Телеграмм: {data['seventh_slide']['telegram']}\n"
        for paragraph in slide.shapes[4].text_frame.paragraphs:
            paragraph.font.size = Pt(16)

    first()
    second()
    third()
    fifth()
    sixth()
    seventh()
    pres.save(f'{input_json["path"]}slides.pptx')